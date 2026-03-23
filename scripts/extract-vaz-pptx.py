from pptx import Presentation
import sys

base = '/Users/diegovillota/Desktop/VAZ Vehículos 2026/Condiciones póliza VH + Beneficios VAZ • 2026/'
files = ['PRESENTACIÓN SIERRA NORTE VH.pptx', 'VEHÍCULOS SIERRA SUR.pptx', 'PRESENTACIÓN COSTA VH.pptx']

for filename in files:
    path = base + filename
    print(f'\n========== {filename} ==========')
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_data:
                            texts.append('TABLE: ' + row_data)
            if texts:
                print(f'-- Slide {i+1} --')
                for t in texts:
                    print(t[:400])
    except Exception as e:
        print(f'Error: {e}')
