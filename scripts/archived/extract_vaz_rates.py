from pptx import Presentation
import os

BASE = "/Users/diegovillota/Desktop/VAZ Vehículos 2026/Condiciones póliza VH + Beneficios VAZ • 2026"
files = [
    os.path.join(BASE, "PRESENTACIÓN SIERRA NORTE VH.pptx"),
    os.path.join(BASE, "VEHÍCULOS SIERRA SUR.pptx"),
    "/tmp/costa_vh.pptx",
]

for f in files:
    print(f"\n=== {os.path.basename(f)} ===")
    prs = Presentation(f)
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))
        if texts:
            print(f"-- Slide {i+1} --")
            for t in texts:
                print(t)
